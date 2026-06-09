#!/usr/bin/env python3
"""PublicLogic — one-page 'at a glance'. Single slide. The whole thing on one page."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
CREAM=RGBColor(0xF5,0xF4,0xF0); SLATE=RGBColor(0x1A,0x1D,0x20); GREEN=RGBColor(0x0F,0x4C,0x3A)
GOLD=RGBColor(0xA9,0x77,0x2F); MUTE=RGBColor(0x6B,0x66,0x60); WHITE=RGBColor(0xFF,0xFF,0xFF)
CREAMW=RGBColor(0xEC,0xEA,0xE3); SERIF="Georgia"; SANS="Helvetica Neue"
EW,EH=Inches(13.333),Inches(7.5)
p=Presentation(); p.slide_width=EW; p.slide_height=EH
s=p.slides.add_slide(p.slide_layouts[6])
def rect(x,y,w,h,c):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); r.shadow.inherit=False; return r
def tx(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.06):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,para in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pp.alignment=align; pp.line_spacing=sp
        for (t,fn,sz,col,bold,ital) in para:
            rn=pp.add_run(); rn.text=t; f=rn.font; f.name=fn; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.italic=ital
rect(0,0,EW,EH,CREAM)
rect(0,0,Inches(0.22),EH,GREEN)
tx(Inches(0.7),Inches(0.45),Inches(11),Inches(0.4),[[("P U B L I C L O G I C   ·   C A P A B I L I T I E S   A T   A   G L A N C E",SANS,11,GOLD,True,False)]])
rect(Inches(0.72),Inches(0.95),Inches(1.0),Pt(2.2),GOLD)
# Hero — North Star
tx(Inches(0.7),Inches(1.2),Inches(12),Inches(1.4),[[("Every system should make it easier for the",SERIF,30,SLATE,True,False)],[("next person to do the right thing.",SERIF,30,GREEN,True,True)]],sp=1.05)
tx(Inches(0.7),Inches(2.75),Inches(12),Inches(0.4),[[("Implementation is the service.  Institutional Stewardship is the method.",SANS,14,GOLD,True,False)]])
# Two columns
tx(Inches(0.7),Inches(3.5),Inches(6),Inches(0.35),[[("WHAT WE DO",SANS,12,GREEN,True,False)]])
svc=[("Project Development","turn ideas into actionable projects"),("Funding Strategy","find and win sustainable funding"),
 ("Implementation Support","move plans into execution"),("Capacity Support","step in and hold the role")]
runs=[]
for h,d in svc: runs.append([("—  ",SANS,12.5,GOLD,True,False),(h+" ",SANS,12.5,SLATE,True,False),("· "+d,SANS,12.5,MUTE,False,False)])
tx(Inches(0.7),Inches(3.9),Inches(6),Inches(2.2),runs,sp=1.35)
rect(Inches(7.1),Inches(3.5),Inches(5.5),Inches(2.55),CREAMW); rect(Inches(7.1),Inches(3.5),Inches(0.1),Inches(2.55),GOLD)
tx(Inches(7.4),Inches(3.7),Inches(5),Inches(0.35),[[("WHY IT'S DIFFERENT",SANS,12,GREEN,True,False)]])
tx(Inches(7.4),Inches(4.1),Inches(5),Inches(1.0),[[("Most systems manage work. PublicLogic helps steward what has to survive the work.",SERIF,15,SLATE,False,True)]],sp=1.1)
tx(Inches(7.4),Inches(5.25),Inches(5),Inches(0.7),[[("We prove it:  ",SANS,11.5,GREEN,True,False),("continuity · record · readiness · adoption · outcome.",SANS,11.5,SLATE,False,False)]],sp=1.1)
# How to start
rect(Inches(0.7),Inches(6.25),Inches(11.9),Inches(0.62),GREEN)
tx(Inches(0.95),Inches(6.25),Inches(11.4),Inches(0.62),[[("HOW TO START   ",SANS,11.5,GOLD,True,False),("Signal → Fit → Map → Build → Sustain → Prove.   Begin with a paid Map.",SANS,12.5,CREAM,True,False)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
tx(Inches(0.7),Inches(7.05),Inches(11.9),Inches(0.3),[[("PublicLogic LLC  ·  nate@publiclogic.org  ·  978-807-0829  ·  Private & Confidential",SANS,8.5,MUTE,False,False)]])
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"final_deck","PublicLogic - One Pager (2026-06).pptx")
p.save(out); print("wrote",out)
