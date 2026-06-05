#!/usr/bin/env python3
"""
Michigan LTC — FINAL institutional deck (14 slides), built to the supplied blueprint.
Palette: cream #F5F4F0 · deep slate #1A1D20 · forest green #0F4C3A · ochre/terracotta accent.
Serif titles (Georgia), sans body (Helvetica Neue). 16:9. Private & Confidential.
Slide 8 disclaimer tightened for internal consistency (recovered-carbon vs $0 pro-forma line).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CREAM=RGBColor(0xF5,0xF4,0xF0); SLATE=RGBColor(0x1A,0x1D,0x20); GREEN=RGBColor(0x0F,0x4C,0x3A)
GOLD=RGBColor(0xA9,0x77,0x2F); MUTE=RGBColor(0x6B,0x66,0x60); WHITE=RGBColor(0xFF,0xFF,0xFF)
CREAMW=RGBColor(0xEC,0xEA,0xE3); GREENL=RGBColor(0xE3,0xEA,0xE5)
SERIF="Georgia"; SANS="Helvetica Neue"
EW,EH=Inches(13.333),Inches(7.5)

def deck():
    p=Presentation(); p.slide_width=EW; p.slide_height=EH; return p
def base(p,bg=CREAM):
    s=p.slides.add_slide(p.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,color,line=None,lw=0.75):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=color
    if line: r.line.color.rgb=line; r.line.width=Pt(lw)
    else: r.line.fill.background()
    r.shadow.inherit=False; return r
def tx(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.04):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pp.alignment=align; pp.line_spacing=sp
        if i>0: pp.space_before=Pt(0)
        for (t,fn,sz,col,bold,ital) in para:
            rn=pp.add_run(); rn.text=t; f=rn.font; f.name=fn; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.italic=ital
    return tb
def kick(s,t,x=Inches(0.85),y=Inches(0.55),col=GOLD): tx(s,x,y,Inches(10),Inches(0.4),[[("   ".join(t.upper()),SANS,11,col,True,False)]])
def rule(s,x=Inches(0.85),y=Inches(1.42),w=Inches(0.9),col=GOLD): rect(s,x,y,w,Pt(2.2),col)
def title(s,t,y=Inches(0.95),size=30,col=SLATE,x=Inches(0.85),w=Inches(11.6)): tx(s,x,y,w,Inches(1.0),[[(t,SERIF,size,col,True,False)]],sp=1.0)
def foot(s,col=MUTE): tx(s,Inches(0.85),Inches(7.06),Inches(11.6),Inches(0.3),[[("MICHIGAN LTC NETWORK   ·   PRIVATE & CONFIDENTIAL   ·   ACCREDITED-INVESTOR DILIGENCE ONLY",SANS,8,col,False,False)]])
def pg(s,n,col=MUTE): tx(s,Inches(12.55),Inches(7.02),Inches(0.6),Inches(0.3),[[(str(n),SANS,9,col,False,False)]],PP_ALIGN.RIGHT)

def section(s,x,y,w,header,bullets,hcol=GREEN,bcol=SLATE,hsize=13,bsize=12.5):
    tx(s,x,y,w,Inches(0.5),[[(header,SANS,hsize,hcol,True,False)]])
    runs=[]
    for b in bullets:
        runs.append([("—  ",SANS,bsize,GOLD,True,False),(b,SANS,bsize,bcol,False,False)])
    tx(s,x,y+Inches(0.52),w,Inches(4.5),runs,sp=1.12)

def chip(s,x,y,w,h,label,fill=GREEN,fg=WHITE):
    rect(s,x,y,w,h,fill); tx(s,x,y,w,h,[[(label,SANS,13,fg,True,False)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)

def arrow(s,x,y,w=Inches(0.5),col=GOLD):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,Inches(0.22)); a.fill.solid(); a.fill.fore_color.rgb=col; a.line.fill.background(); a.shadow.inherit=False

OUT=os.path.join(os.path.dirname(os.path.abspath(__file__)),"final_deck")
os.makedirs(OUT,exist_ok=True)
p=deck()

# S1 — TITLE
s=base(p,SLATE)
rect(s,Inches(11.45),0,Inches(1.88),EH,GREEN)
tx(s,Inches(11.45),Inches(3.2),Inches(1.88),Inches(1),[[("upgrade",SERIF,14,RGBColor(0xCF,0xC9,0xBE),False,True)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
kick(s,"Platform Position",col=GOLD)
rule(s,y=Inches(1.35),w=Inches(1.1))
tx(s,Inches(0.85),Inches(1.9),Inches(10.3),Inches(2.0),[[("MICHIGAN LTC NETWORK",SERIF,46,CREAM,True,False)],[("A Circular-Industrial Platform",SERIF,22,GOLD,False,True)]],sp=1.05)
tx(s,Inches(0.85),Inches(4.25),Inches(10.0),Inches(0.9),[[("“We don’t process waste. We upgrade stranded material.”",SERIF,20,CREAM,False,True)]],sp=1.05)
pills=["Carbon","Capital","Community","Continuity"]
for i,pl in enumerate(pills): chip(s,Inches(0.85)+Inches(2.55)*i,Inches(5.4),Inches(2.35),Inches(0.55),pl,GREEN,CREAM)
tx(s,Inches(0.85),Inches(6.25),Inches(10.4),Inches(1.0),[
 [("Private & Confidential",SANS,11,GOLD,True,False)],
 [("For accredited-investor diligence only. Subject to definitive documentation and securities-counsel review. Base figures reconciled to a single source of truth; RIN/§45Z/SAF credits are contingent and excluded from the base case. Technology used under license.",SANS,8.5,RGBColor(0xB9,0xB3,0xA8),False,False)]
],sp=1.05)

# S2 — STRATEGIC PIVOT
s=base(p); kick(s,"The Strategic Pivot"); rule(s); title(s,"Asset Conversion over Waste Management")
cols=[(Inches(0.85),"The Massive Regional Liability",["Millions of tons of tires, treated wood, and biomass sit stranded in regional supply chains.","Traditional recycling breaks down under the volume and negative unit economics."]),
 (Inches(4.78),"The Material Upgrade Pivot",["We co-locate directly with the industrial hosts who already control the material.","We convert raw physical liabilities into higher-value energy and recovered carbon."]),
 (Inches(8.71),"Commercial Impact",["Shorter path to revenue execution.","Radically lowered host infrastructure risk.","A repeatable state-by-state model under one common governance spine."])]
for x,h,b in cols: section(s,x,Inches(2.0),Inches(3.7),h,b)
foot(s); pg(s,2)

# S3 — RESILIENT SYSTEM
s=base(p); kick(s,"Building a Resilient System"); rule(s); title(s,"A Networked Platform, Not a Single Asset")
items=[("01","Partner, Don’t Compete",["Site alongside operators who already aggregate feedstock.","Create an immediate, profitable new revenue line for them.","Secure accelerated site control, logistics, and permitting."]),
 ("02","A Network, Not a Plant",["A state-by-state model built on pre-existing industrial infrastructure.","De-risks each subsequent site and compresses time-to-market."]),
 ("03","Governance as the Spine",["Risk ring-fenced and isolated at the site level via dedicated SPVs.","Intelligence, compliance, and replication scale at the parent level."])]
for i,(num,h,b) in enumerate(items):
    x=Inches(0.85)+Inches(3.95)*i
    tx(s,x,Inches(1.95),Inches(1.2),Inches(0.8),[[(num,SERIF,30,GOLD,True,False)]])
    section(s,x,Inches(2.75),Inches(3.7),h,b)
foot(s); pg(s,3)

# S4 — CORRIDOR
s=base(p); kick(s,"The Michigan Corridor"); rule(s); title(s,"Three Active Nodes. One Repeatable Model.")
nodes=[("Flint / Dort Hwy","EJ Anchor Node",GREEN),("Coleman / Midland","Replication Template",GREEN),("Coldwater Road","Phase 2 Infrastructure",SLATE)]
bx=Inches(0.85); bw=Inches(3.5); gap=Inches(0.62); by=Inches(1.95); bh=Inches(1.0)
for i,(n,sub,col) in enumerate(nodes):
    x=bx+(bw+gap)*i
    rect(s,x,by,bw,bh,col)
    tx(s,x,by+Inches(0.14),bw,Inches(0.5),[[(n,SANS,14,CREAM,True,False)]],PP_ALIGN.CENTER)
    tx(s,x,by+Inches(0.56),bw,Inches(0.4),[[(sub,SANS,10.5,RGBColor(0xCF,0xC9,0xBE),False,True)]],PP_ALIGN.CENTER)
    if i<2: arrow(s,x+bw+Inches(0.06),by+Inches(0.4),Inches(0.5))
desc=[("Node 01 · Flint / Dort Highway",["14-acre brownfield co-located with Environmental Rubber Recycling (ERR).","~40,000 tons/yr qualified waste-tire supply under a 10-year arrangement.","The corridor’s environmental-justice and brownfield anchor."]),
 ("Node 02 · Coleman / Midland",["Rural co-location with Geocycle, an active Part 169-registered facility.","~22,000 tons/yr existing tire-derived throughput.","The template for co-location and USDA-aligned rural financing."]),
 ("Node 03 · Coldwater Road",["118.3-acre rail-served hub; 4,570 ft of Lake State Railway frontage.","Heavy industrial (I-2); EGLE No Further Action complete with controls.","Held for Phase 2 large-scale logistics and regional scaling."])]
for i,(h,b) in enumerate(desc):
    section(s,bx+(bw+gap)*i,Inches(3.35),Inches(3.6),h,b,bsize=11)
foot(s); pg(s,4)

# S5 — WHY MICHIGAN
s=base(p); kick(s,"Why Michigan Wins"); rule(s); title(s,"Geography and Structural Moats")
grid=[("Unmatched Feedstock Density",["High regional concentration of automotive scrap, industrial byproducts, and agricultural / forestry residues."]),
 ("Logistics and Infrastructure",["Pre-existing heavy rail, deep utilities, and a deep inventory of underutilized industrial brownfields."]),
 ("Industrial Labor Base",["Immediate access to a manufacturing-ready workforce specialized in 24/7 process mechanics."]),
 ("Ecosystem Synergy",["Multi-year relationships with dominant regional recyclers and hosts.","Active state support for circular transitions (NextCycle, EGLE)."])]
for i,(h,b) in enumerate(grid):
    x=Inches(0.85)+Inches(6.1)*(i%2); y=Inches(2.0)+Inches(2.45)*(i//2)
    section(s,x,y,Inches(5.6),h,b)
foot(s); pg(s,5)

# S6 — OPERATIONAL FLOW
s=base(p); kick(s,"The Operational Flow"); rule(s); title(s,"Controlled Intake to Local Value Creation")
steps=[("01","Source","Qualify stream & custody (host assets)"),("02","Prepare","Sort, size, and stage feedstocks"),
 ("03","Convert","Advanced pyrolysis: clean carbon & gas"),("04","Recover","Fuel + carbon product streams out"),
 ("05","Validate","Test, certify, and reconcile"),("06","Reinvest","Sell, report, and return value locally")]
bw=Inches(1.86); bh=Inches(2.2); by=Inches(2.2); x0=Inches(0.7); g=Inches(0.13)
for i,(num,h,d) in enumerate(steps):
    x=x0+(bw+g)*i
    rect(s,x,by,bw,bh,CREAMW if i%2==0 else GREENL)
    rect(s,x,by,bw,Inches(0.12),GOLD)
    tx(s,x+Inches(0.15),by+Inches(0.28),bw-Inches(0.3),Inches(0.6),[[(num,SERIF,26,GOLD,True,False)]])
    tx(s,x+Inches(0.15),by+Inches(0.95),bw-Inches(0.3),Inches(0.4),[[(h,SANS,13,SLATE,True,False)]])
    tx(s,x+Inches(0.15),by+Inches(1.35),bw-Inches(0.3),Inches(0.8),[[(d,SANS,9.5,MUTE,False,False)]],sp=1.05)
tx(s,Inches(0.85),Inches(5.05),Inches(11.6),Inches(0.9),[
 [("The operating standard",SANS,13,GREEN,True,False)],
 [("Proceed with deployment only when feedstock stream, conversion technology, compliance evidence, target markets, and local community value are completely legible.",SANS,12.5,SLATE,False,False)]],sp=1.1)
foot(s); pg(s,6)

# S7 — GUARDRAILS
s=base(p); kick(s,"Technical & Operational Guardrails"); rule(s); title(s,"Proven Equipment. Honest Labor Realities.")
g=[("Licensed Technology Deployment",["Proprietary WasteWerx depolymerization systems under strict operating license.","Contained, safe process loops with favorable environmental discharge."]),
 ("Isolation of R&D Risk",["The platform is insulated from fundamental technology-development risk.","WasteWerx retains the IP; Michigan LTC monetizes the operating rights."]),
 ("Operator-Reality Staffing",["We reject unrealistic automated “10-person” assumptions for heavy infrastructure.","A fully-costed 62-FTE local workforce per site for true, safe 24/7 uptime."])]
for i,(h,b) in enumerate(g):
    section(s,Inches(0.85)+Inches(3.95)*i,Inches(2.05),Inches(3.7),h,b)
foot(s); pg(s,7)

# S8 — ECONOMICS
s=base(p); kick(s,"Reconciled Site Economics"); rule(s); title(s,"We Bank What We Can Support")
# two metric cards
cards=[("Modeled Annual Revenue","$16.88M","Baseline of contractable physical commodities (fuel + recovered carbon), valued at conservative market."),
 ("Modeled Annual EBITDA","~$6.7M","Reconciled against the operator-reality 62-FTE staffing floor and feedstock opex.")]
for i,(lab,num,note) in enumerate(cards):
    x=Inches(0.85)+Inches(6.0)*i
    rect(s,x,Inches(2.0),Inches(5.6),Inches(2.0),CREAMW); rect(s,x,Inches(2.0),Inches(0.12),Inches(2.0),GREEN)
    tx(s,x+Inches(0.35),Inches(2.2),Inches(5.0),Inches(0.4),[[(lab.upper(),SANS,11,GREEN,True,False)]])
    tx(s,x+Inches(0.35),Inches(2.55),Inches(5.0),Inches(0.9),[[(num,SERIF,48,SLATE,True,False)]])
    tx(s,x+Inches(0.35),Inches(3.5),Inches(5.0),Inches(0.5),[[(note,SANS,10.5,MUTE,False,False)]],sp=1.05)
tx(s,Inches(0.85),Inches(4.45),Inches(11.6),Inches(0.4),[[("Strict underwriting disclaimers",SANS,13,GREEN,True,False)]])
disc=[("Recovered carbon:  ","the executed WasteWerx Pro Forma prices fuel only (carbon line at $0). The base values recovered carbon at conservative market and must be confirmed against a carbon offtake — treat as uncontracted until then."),
 ("Environmental incentives:  ","all regulatory credits (RFS / D4 / D7 RINs, §45Z clean-fuel) are segregated from base-case returns and treated strictly as contingent, post-qualification upside (~$20.4M/yr).")]
runs=[]
for head,body in disc: runs.append([("—  ",SANS,12,GOLD,True,False),(head,SANS,12,SLATE,True,False),(body,SANS,12,SLATE,False,False)])
tx(s,Inches(0.85),Inches(4.95),Inches(11.6),Inches(1.8),runs,sp=1.14)
foot(s); pg(s,8)

# S9 — WHAT EXISTS TODAY
s=base(p); kick(s,"What Exists Today"); rule(s); title(s,"Diligence Readiness and Execution Baseline")
section(s,Inches(0.85),Inches(2.0),Inches(5.7),"Current Asset Posture",[
 "Site control: co-location host arrangements established with ERR and Geocycle.",
 "Feedstock footprint: ~62,000 tons/yr identified and documented across active nodes.",
 "Environmental: EGLE No Further Action complete with controls on key corridor hubs.",
 "Replication baseline: moved from legacy assumptions to line-by-line engineering budgets."])
section(s,Inches(6.95),Inches(2.0),Inches(5.5),"Remaining Risks & Mitigants",[
 "Permit sequencing — mitigated by co-locating on pre-permitted, industrial-zoned sites.",
 "Deployment schedules — mitigated by modular, factory-tested licensed equipment.",
 "Offtake finalization — advancing documented fuel offtake MOUs before commissioning."])
foot(s); pg(s,9)

# S10 — GOVERNANCE
s=base(p); kick(s,"PublicLogic Governance Infrastructure"); rule(s); title(s,"The Platform Integration Layer")
rect(s,Inches(4.4),Inches(2.0),Inches(4.5),Inches(0.7),GREEN)
tx(s,Inches(4.4),Inches(2.0),Inches(4.5),Inches(0.7),[[("PUBLICLOGIC COMPLIANCE FRAMEWORK",SANS,12,CREAM,True,False)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
for i,(h,items) in enumerate([("Governance Architecture",["State SPV ring-fencing","Replication playbooks","Continuity controls"]),
                              ("Evidence Management",["Audit-ready records","Cost allowability","Grant-stack integration"])]):
    x=Inches(1.7)+Inches(5.6)*i
    rect(s,x,Inches(3.1),Inches(4.4),Inches(2.0),CREAMW); rect(s,x,Inches(3.1),Inches(4.4),Inches(0.1),GOLD)
    tx(s,x+Inches(0.3),Inches(3.3),Inches(3.8),Inches(0.4),[[(h,SANS,13,GREEN,True,False)]])
    tx(s,x+Inches(0.3),Inches(3.8),Inches(3.8),Inches(1.2),[[("—  ",SANS,12,GOLD,True,False),(it,SANS,12,SLATE,False,False)] for it in items],sp=1.2)
tx(s,Inches(0.85),Inches(5.4),Inches(11.6),Inches(1.2),[
 [("PublicLogic provides the institutional spine — governance, documentation, non-dilutive grant-stack coordination, and compliance continuity. The platform scales state-by-state because the operating system is built on standardized, auditable evidence. PublicLogic is the governance partner, not a placement agent; its fee is fixed and non-contingent.",SANS,12,SLATE,False,False)]],sp=1.12)
foot(s); pg(s,10)

# S11 — COMMUNITY
s=base(p); kick(s,"Community-Embedded Design"); rule(s); title(s,"The Return Circulates Locally")
section(s,Inches(0.85),Inches(2.0),Inches(11.6),"The Structural Commitment",[
 "A contractual 10% share of net profit at each operating site is directed to locally governed community development — encoded in the operating agreements, not treated as afterthought philanthropy."])
two=[("The Warehouse / Gen 215 · Flint Anchor",["Secure community infrastructure, family support networks, and localized urban food production."]),
 ("ACS Cooperative Partnership",["High-density aquaponic programming and tech-workforce training linked to university research."])]
for i,(h,b) in enumerate(two): section(s,Inches(0.85)+Inches(6.0)*i,Inches(3.6),Inches(5.6),h,b)
tx(s,Inches(0.85),Inches(5.6),Inches(11.6),Inches(0.6),[[("Note: community-benefit basis (% of net profit vs % of equity) and any housing/PUD mechanism are pending reconciliation before publication.",SANS,10,MUTE,False,True)]])
foot(s); pg(s,11)

# S12 — THE ASK
s=base(p); kick(s,"The Institutional Ask"); rule(s); title(s,"Two Tranches. One Disciplined Capital Stack.")
asks=[("Project Deployment","$45,000,000","$15,000,000 per site × 3","Site prep, core hardware, integration, and site-level commissioning.",GREEN),
 ("Corporate Platform","$2,500,000","Parent holdco equity","Governance, replication/IP architecture, multi-state BD, diligence & legal readiness. Allocations set before definitive docs.",SLATE),
 ("Grant-Stack Coordination","$4M–$14M","Non-dilutive offsets","NextCycle MI, USDA REAP, EPA Brownfield lanes — de-risking early private equity. Distinct from the securities offering.",GOLD)]
for i,(h,num,sub,note,col) in enumerate(asks):
    x=Inches(0.85)+Inches(4.05)*i
    rect(s,x,Inches(2.0),Inches(3.8),Inches(3.4),CREAMW); rect(s,x,Inches(2.0),Inches(3.8),Inches(0.13),col)
    tx(s,x+Inches(0.3),Inches(2.25),Inches(3.3),Inches(0.4),[[(h.upper(),SANS,11,col if col!=SLATE else GREEN,True,False)]])
    tx(s,x+Inches(0.3),Inches(2.65),Inches(3.3),Inches(0.8),[[(num,SERIF,30,SLATE,True,False)]])
    tx(s,x+Inches(0.3),Inches(3.5),Inches(3.3),Inches(0.4),[[(sub,SANS,11.5,GOLD,True,False)]])
    tx(s,x+Inches(0.3),Inches(3.95),Inches(3.3),Inches(1.3),[[(note,SANS,11,MUTE,False,False)]],sp=1.1)
tx(s,Inches(0.85),Inches(5.7),Inches(11.6),Inches(0.5),[[("Company raise is presented as a distinct ask and is never blended into site totals. Use-of-funds and return structure subject to securities-counsel review.",SANS,10,MUTE,False,True)]])
foot(s); pg(s,12)

# S13 — ENDURANCE THESIS
s=base(p,GREEN)
rule(s,y=Inches(2.4),w=Inches(1.1),col=GOLD)
tx(s,Inches(0.85),Inches(1.6),Inches(11.6),Inches(0.6),[[("The Governance Is The Product.",SERIF,34,CREAM,True,False)]])
tx(s,Inches(0.85),Inches(2.8),Inches(11.2),Inches(2.2),[[("“We close like an institutional process. We distinguish — out loud — between what is real, what is advanced, and what is still being locked. The base case is supported by modeled throughput. The upside is highly disciplined.”",SERIF,21,CREAM,False,True)]],sp=1.12)
tx(s,Inches(0.85),Inches(5.5),Inches(11.6),Inches(1),[[("MICHIGAN LTC NETWORK",SERIF,22,WHITE,True,False)],[("Carbon  ·  Capital  ·  Community  ·  Continuity",SANS,13,GOLD,True,False)]],sp=1.3)

# S14 — APPENDIX
s=base(p); kick(s,"Appendix — Basis of Record"); rule(s); title(s,"Data Room Line-for-Line Traceability")
tx(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(0.5),[[("“Every figure is a signed-document value or a formula over one.”",SERIF,15,GREEN,False,True)]])
rowsT=[["Core Parameter","Documented Source Baseline","Source ID"],
 ["Revenue projections","WasteWerx Model 3 License Pro Forma P&L","SRC-003"],
 ["Capital budgets","Site development budget & vendor cost build-up","Investor Workbook v2.0"],
 ["Staffing & local opex","Operator-reality 62-FTE site staffing model","SRC-001"],
 ["Grant underwriting","Published agency program guidance (EGLE / USDA)","SRC-002"],
 ["Governance architecture","PublicLogic cyclical-services control framework","PL-Model-2026"]]
gt=s.shapes.add_table(len(rowsT),3,Inches(0.85),Inches(2.6),Inches(11.6),Inches(3.3)).table
gt.columns[0].width=Inches(3.4); gt.columns[1].width=Inches(5.7); gt.columns[2].width=Inches(2.5)
for j,h in enumerate(rowsT[0]):
    c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=SLATE
    r=c.text_frame.paragraphs[0].add_run(); r.text=h; r.font.name=SANS; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=CREAM
for i,row in enumerate(rowsT[1:]):
    for j,v in enumerate(row):
        c=gt.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=CREAM if i%2==0 else CREAMW
        r=c.text_frame.paragraphs[0].add_run(); r.text=v; r.font.name=SANS; r.font.size=Pt(12); r.font.color.rgb=SLATE
        if j==2: r.font.bold=True; r.font.color.rgb=GREEN
foot(s); pg(s,14)

out=os.path.join(OUT,"Michigan LTC - FINAL Deck (institutional 2026-06-05).pptx")
p.save(out); print("wrote",out,"| slides:",len(p.slides.__iter__.__self__._sldIdLst))
