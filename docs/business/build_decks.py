#!/usr/bin/env python3
"""
Builds three elevated, 'coy' decks for the Michigan LTC corridor:
  A — Institutional Memory (PublicLogic, rethought)
  C — Impact (the lead)
  B — Money (restrained / withholding)
Design: deep ink + warm paper, restrained serif display, generous whitespace,
art-directed photographic placeholders (alternating light / duotone wells).
Financials are unchanged from Canonical_Financial_Model v2.0. Private & Confidential.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

INK=RGBColor(0x17,0x20,0x2B); PAPER=RGBColor(0xF6,0xF3,0xEE); GOLD=RGBColor(0xB0,0x89,0x48)
WELL_L=RGBColor(0xE7,0xE1,0xD7); WELL_D=RGBColor(0x23,0x2E,0x38); MUTE=RGBColor(0x7A,0x71,0x66)
PAPER_ON_INK=RGBColor(0xEC,0xE7,0xDE); SERIF="Georgia"; SANS="Arial"
EW,EH=Inches(13.333),Inches(7.5)

def deck():
    p=Presentation(); p.slide_width=EW; p.slide_height=EH; return p
def slide(p,bg=PAPER):
    s=p.slides.add_slide(p.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,color,line=None):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=color
    if line: r.line.color.rgb=line; r.line.width=Pt(0.75)
    else: r.line.fill.background()
    r.shadow.inherit=False; return r
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.04):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pp.alignment=align; pp.line_spacing=space
        for (t,fn,sz,col,bold,ital,track) in para:
            rn=pp.add_run(); rn.text=t; f=rn.font; f.name=fn; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.italic=ital
    return tb
def kick(s,x,y,t,col=GOLD): txt(s,x,y,Inches(8),Inches(0.4),[[(t.upper(),SANS,11.5,col,True,False,0)]])
def rule(s,x,y,w=Inches(0.9),col=GOLD): rect(s,x,y,w,Pt(2),col)
def well(s,x,y,w,h,dark,subject):
    rect(s,x,y,w,h,WELL_D if dark else WELL_L)
    cap=f"[ photograph — {subject} ]"
    txt(s,x+Inches(0.3),y,w-Inches(0.6),h,[[(cap,SERIF,12.5,(GOLD if dark else MUTE),False,True,0)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
def page(s,n,col=MUTE): txt(s,Inches(12.1),Inches(7.0),Inches(1.0),Inches(0.3),[[(n,SANS,9,col,False,False,0)]],PP_ALIGN.RIGHT)
def confid(s,col=MUTE): txt(s,Inches(0.7),Inches(7.0),Inches(6),Inches(0.3),[[("PRIVATE & CONFIDENTIAL",SANS,8.5,col,True,False,0)]])

def cover(p,kicker,lines,sub):
    s=slide(p,INK); rule(s,Inches(0.9),Inches(1.5),Inches(1.1))
    kick(s,Inches(0.9),Inches(1.0),kicker)
    runs=[[(ln,SERIF,46,PAPER_ON_INK,False,it,0)] for ln,it in lines]
    txt(s,Inches(0.9),Inches(2.1),Inches(11.5),Inches(3.2),runs,space=1.06)
    txt(s,Inches(0.9),Inches(6.3),Inches(11),Inches(0.6),[[(sub,SANS,12.5,GOLD,False,False,0)]])
    # right-edge duotone band for richness
    rect(s,Inches(11.4),0,Inches(1.93),EH,WELL_D)
    txt(s,Inches(11.4),Inches(3.2),Inches(1.93),Inches(1),[[("the record",SERIF,13,GOLD,False,True,0)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    return s

def statement(p,lines,sub=None,n=None):
    s=slide(p); rule(s,Inches(0.9),Inches(1.4))
    runs=[[(t,SERIF,40,(GOLD if it else INK),False,it,0)] for t,it in lines]
    txt(s,Inches(0.9),Inches(2.0),Inches(10.5),Inches(3.2),runs,space=1.08)
    if sub: txt(s,Inches(0.9),Inches(5.6),Inches(10),Inches(0.8),[[(sub,SANS,13,MUTE,False,False,0)]])
    if n: page(s,n)
    return s

def split(p,kicker,lines,subject,dark,sub=None,n=None,side="right"):
    s=slide(p)
    if side=="right":
        well(s,Inches(7.6),0,Inches(5.733),EH,dark,subject); tx=Inches(0.9); tw=Inches(6.2)
    else:
        well(s,0,0,Inches(5.733),EH,dark,subject); tx=Inches(6.4); tw=Inches(6.2)
    rule(s,tx,Inches(1.4)); kick(s,tx,Inches(1.0),kicker)
    runs=[[(t,SERIF,33,(GOLD if it else INK),False,it,0)] for t,it in lines]
    txt(s,tx,Inches(2.0),tw,Inches(3.4),runs,space=1.08,anchor=MSO_ANCHOR.TOP)
    if sub: txt(s,tx,Inches(5.4),tw,Inches(1.4),[[(sub,SANS,12.5,MUTE,False,False,0)]])
    if n: page(s,n)
    return s

def fullimage(p,subject,overlay,sub=None,n=None):
    s=slide(p,INK); well(s,0,0,EW,EH,True,subject)
    rect(s,0,Inches(4.6),EW,Inches(2.9),INK)  # lower scrim
    rule(s,Inches(0.9),Inches(5.0))
    txt(s,Inches(0.9),Inches(5.2),Inches(11),Inches(1.4),[[(t,SERIF,34,PAPER_ON_INK,False,it,0)] for t,it in overlay],space=1.06)
    if sub: txt(s,Inches(0.9),Inches(6.6),Inches(11),Inches(0.5),[[(sub,SANS,12,GOLD,False,False,0)]])
    if n: page(s,n,PAPER_ON_INK)
    return s

def places(p,title,items,n=None):  # items: [(label, subject, dark)]
    s=slide(p); rule(s,Inches(0.9),Inches(0.95)); kick(s,Inches(0.9),Inches(0.6),"THE CORRIDOR, IN REAL PLACES")
    txt(s,Inches(0.9),Inches(1.15),Inches(11),Inches(0.9),[[(title,SERIF,30,INK,False,False,0)]])
    x=Inches(0.9); w=Inches(3.71); gap=Inches(0.3); y=Inches(2.4); h=Inches(3.6)
    for i,(lab,subj,dark) in enumerate(items):
        xx=x+(w+gap)*i
        well(s,xx,y,w,h,dark,subj)
        txt(s,xx,y+h+Inches(0.15),w,Inches(0.4),[[(lab,SANS,13,INK,True,False,0)]],PP_ALIGN.CENTER)
    if n: page(s,n)
    return s

def bignum(p,num,label,withhold,n=None):
    s=slide(p,INK); rule(s,Inches(0.9),Inches(1.5),Inches(1.1))
    kick(s,Inches(0.9),Inches(1.05),"PER SITE")
    txt(s,Inches(0.8),Inches(1.9),Inches(11.5),Inches(3.0),[[(num,SERIF,150,PAPER_ON_INK,False,False,0)]],space=0.9)
    txt(s,Inches(0.95),Inches(4.7),Inches(8),Inches(0.6),[[(label,SANS,15,GOLD,True,False,0)]])
    txt(s,Inches(0.95),Inches(5.6),Inches(10.5),Inches(1.2),[[(withhold,SERIF,20,PAPER_ON_INK,False,True,0)]],space=1.08)
    confid(s,RGBColor(0x9A,0x90,0x83));
    if n: page(s,n,PAPER_ON_INK)
    return s

def table(p,title,headers,rows,caption,n=None):
    s=slide(p); rule(s,Inches(0.9),Inches(0.95)); kick(s,Inches(0.9),Inches(0.6),"PER-SITE ECONOMICS — DILIGENCE VIEW")
    txt(s,Inches(0.9),Inches(1.15),Inches(11),Inches(0.8),[[(title,SERIF,28,INK,False,False,0)]])
    nrows=len(rows)+1; tw=Inches(11.5); th=Inches(0.55)*nrows
    gt=s.shapes.add_table(nrows,len(headers),Inches(0.9),Inches(2.2),tw,th).table
    gt.columns[0].width=Inches(6.5)
    for j in range(1,len(headers)): gt.columns[j].width=Inches(int((11.5-6.5)*914400/ (len(headers)-1)))
    for j,h in enumerate(headers):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=INK
        para=c.text_frame.paragraphs[0]; r=para.add_run(); r.text=h; r.font.name=SANS; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=PAPER_ON_INK
        para.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=gt.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=PAPER if i%2==0 else WELL_L
            para=c.text_frame.paragraphs[0]; r=para.add_run(); r.text=val; r.font.name=SANS; r.font.size=Pt(11.5); r.font.color.rgb=INK
            para.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
    txt(s,Inches(0.9),Inches(6.5),Inches(11.5),Inches(0.8),[[(caption,SANS,10.5,MUTE,False,True,0)]],space=1.05)
    if n: page(s,n)
    return s

def closer(p,line,sub):
    s=slide(p,INK); rule(s,Inches(0.9),Inches(2.7),Inches(1.1))
    txt(s,Inches(0.9),Inches(3.0),Inches(11),Inches(2),[[(t,SERIF,40,(GOLD if it else PAPER_ON_INK),False,it,0)] for t,it in line],space=1.06)
    txt(s,Inches(0.9),Inches(5.6),Inches(11),Inches(0.5),[[(sub,SANS,12,RGBColor(0x9A,0x90,0x83),False,False,0)]])
    return s

def contact(p,title):
    s=slide(p,INK); rule(s,Inches(0.9),Inches(1.2),Inches(1.1)); kick(s,Inches(0.9),Inches(0.8),title)
    txt(s,Inches(0.9),Inches(2.0),Inches(11),Inches(2),[
        [("Michigan LTC Network",SERIF,30,PAPER_ON_INK,False,False,0)],
        [("Resource Recovery Corridor",SERIF,20,GOLD,False,True,0)],
    ],space=1.1)
    txt(s,Inches(0.9),Inches(4.2),Inches(11),Inches(2),[
        [("Energy Mann & Sunn, Inc.  ·  Robert C. McCall Jr., Project Executive",SANS,12.5,PAPER_ON_INK,False,False,0)],
        [("Governance & Compliance — PublicLogic LLC  ·  Nathan Boudreau, MPA, MCPPO  ·  Dr. Allison Weiss Rothschild",SANS,12.5,PAPER_ON_INK,False,False,0)],
        [("nate@publiclogic.org  ·  978-807-0829",SANS,12.5,GOLD,False,False,0)],
    ],space=1.5)
    confid(s,RGBColor(0x9A,0x90,0x83)); return s

OUT=os.path.join(os.path.dirname(os.path.abspath(__file__)),"casespace","09_DELIVERABLES","06_Investor_or_Partner")
os.makedirs(OUT,exist_ok=True)

# ============ A — INSTITUTIONAL MEMORY ============
A=deck()
cover(A,"PublicLogic",[("The corridor remembers.",False),("So nothing has to be relearned.",True)],"Institutional memory as infrastructure  ·  Private & Confidential")
split(A,"The thesis",[("Most projects forget.",False),("This one remembers.",True)],"a real record — a logbook, a permit wall, a hand resting on a ledger",True,"Knowledge usually lives in people. When they leave, it leaves with them.","02")
split(A,"Memory, made portable",[("Every permit. Every reading.",False),("Every relationship. Every lesson.",False),("Retained.",True)],"close-up of an open field binder, tabs and tabs",False,"The corridor keeps its own record — and carries it forward.","03",side="left")
statement(A,[("Each site starts where",False),("the last one finished.",True)],"Continuity is the compounding engine: the second site is cheaper, faster, and more certain than the first.","04")
split(A,"The evidence layer",[("PuddleJumper is the quiet",False),("memory underneath the work.",False)],"a single screen in a dim control room, a calm interface",True,"Every submission, approval, and deadline — logged, sealed, auditable.","05")
statement(A,[("Proof, not promises.",False)],"Documentation-first. Verification-ready. The record is the product.","06")
closer(A,[("Institutional memory is the",False),("quietest moat there is.",True)],"PublicLogic LLC  ·  Governance · Compliance · Continuity")
contact(A,"PublicLogic — Governance & Continuity")
A.save(os.path.join(OUT,"2026-06-05 - Michigan LTC - A - Institutional Memory - Review Draft.pptx"))

# ============ C — IMPACT (the lead) ============
C=deck()
cover(C,"Michigan LTC Network",[("An everlasting",False),("cycle of good.",True)],"Resource Recovery Corridor  ·  Coleman – Flint  ·  Private & Confidential")
split(C,"The belief",[("Where others see waste,",False),("we see what comes next.",True)],"before / after — a tired industrial yard, then the same yard alive",True,"End-of-life tires, railroad ties, and biomass become fuel, recovered material, and work.","02")
statement(C,[("We don't arrive.",False),("We co-locate.",True)],"Geocycle (Coleman) · Environmental Rubber Recycling (Flint) · National Energy of Lincoln — existing hosts, new value.","03")
places(C,"Real towns. Real hosts. Real jobs.",[
    ("Lincoln · Alcona County","an 18 MW biomass plant, steam and steel, rural Michigan",True),
    ("Flint · Genesee County","Dort Highway — a working tire yard in an EJ community",False),
    ("Coleman · Midland County","Geocycle's recycling line, conveyors and bales",True),
],"04")
statement(C,[("Waste becomes fuel.",False),("Fuel becomes revenue.",False),("Revenue becomes community.",True)],"The economic generator and the community investment are the same system.","05")
split(C,"Community, structurally embedded",[("Ten percent, returned —",False),("by design, not as charity.",True)],"a commercial kitchen incubator, hands and steam, people learning",False,"The Warehouse / Gen 215 · ACS Cooperative — workforce, food security, participation.","06",side="left")
statement(C,[("What it leaves behind:",False),("jobs kept, capacity built,",False),("a record that endures.",True)],"Built to survive turnover, partner changes, and the next twenty-five years.","07")
closer(C,[("Legacy is what the next",False),("generation can use.",True)],"Michigan LTC Network  ·  Resource Recovery Corridor")
contact(C,"Michigan LTC — Partnership")
C.save(os.path.join(OUT,"2026-06-05 - Michigan LTC - C - Impact - Review Draft.pptx"))

# ============ B — MONEY (coy) ============
B=deck()
cover(B,"In Confidence",[("The economics,",False),("withheld until earned.",True)],"For accredited discussion only  ·  Private & Confidential")
statement(B,[("Capital is the seed.",False),("Stewardship is the water.",True)],"Returns come from a system that holds together for twenty-five years — not from a single number.","02")
bignum(B,"$15M","per site  ·  three sites in the corridor","The rest is in the diligence room.","03")
statement(B,[("What we share,",False),("and when.",True)],"Base case openly · contingent upside clearly labeled · the full model on request, under NDA, with counsel.","04")
table(B,"One site, on the executed Pro Forma",
    ["Line","As signed","Operator-reality"],
    [
     ["Year 1 cash out","$7,622,828","$7,622,828"],
     ["Steady-state revenue","$35,581,680","$35,581,680"],
     ["Steady-state EBITDA","$30,200,531","~$22,047,043"],
     ["of which RIN + §45Z (contingent)","$20,428,560","$20,428,560"],
     ["Base fuel only (no credits)","$9,771,971","~$1,618,483"],
    ],
    "Keyed to Canonical v2.0. 'As signed' assumes Pro Forma staffing; operator-reality reflects full staffing + feedstock. "
    "$20.4M of revenue is contingent on RFS/§45Z registration. Carbon black is $0 in the base. Not an offer of securities.","05")
statement(B,[("The corridor, in three.",False)],"$45M across three sites · a separate $2.5M at the company · $47.5M total — presented as distinct asks.","06")
split(B,"What stays private",[("The full model travels",False),("only as far as trust does.",True)],"a closed folder on a desk, a single lamp, nothing on screen",True,"Per-site and network detail, sensitivity, and structure — shared in the room, not on the page.","07",side="left")
closer(B,[("Numbers earn their",False),("audience.",True)],"Private & Confidential · Not an offer to sell securities · Subject to definitive docs + securities counsel")
contact(B,"Michigan LTC — In Confidence")
B.save(os.path.join(OUT,"2026-06-05 - Michigan LTC - B - Money - Confidential Draft.pptx"))

print("wrote 3 decks to",OUT)
for f in sorted(os.listdir(OUT)):
    if f.endswith(".pptx"): print("  ",f)
