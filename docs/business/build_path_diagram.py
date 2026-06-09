#!/usr/bin/env python3
"""PublicLogic — The Path. One-page visual. The diagram that appears everywhere."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
CREAM=RGBColor(0xF5,0xF4,0xF0); SLATE=RGBColor(0x1A,0x1D,0x20); GREEN=RGBColor(0x0F,0x4C,0x3A)
GOLD=RGBColor(0xA9,0x77,0x2F); MUTE=RGBColor(0x6B,0x66,0x60); WHITE=RGBColor(0xFF,0xFF,0xFF)
CREAMW=RGBColor(0xEC,0xEA,0xE3); SERIF="Georgia"; SANS="Helvetica Neue"
EW,EH=Inches(13.333),Inches(7.5)
p=Presentation(); p.slide_width=EW; p.slide_height=EH
s=p.slides.add_slide(p.slide_layouts[6])
def rect(x,y,w,h,c,shape=MSO_SHAPE.RECTANGLE,line=None):
    r=s.shapes.add_shape(shape,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=c
    if line is None: r.line.fill.background()
    else: r.line.color.rgb=line; r.line.width=Pt(0.75)
    r.shadow.inherit=False; return r
def tx(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,sp=1.04):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,para in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pp.alignment=align; pp.line_spacing=sp
        for (t,fn,sz,col,bold,ital) in para:
            rn=pp.add_run(); rn.text=t; f=rn.font; f.name=fn; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.italic=ital
    return tb
rect(0,0,EW,EH,CREAM)
rect(0,0,Inches(0.22),EH,GREEN)
# header
tx(Inches(0.7),Inches(0.4),Inches(12),Inches(0.35),[[("P U B L I C L O G I C   ·   T H E   P A T H",SANS,11,GOLD,True,False)]],anchor=MSO_ANCHOR.TOP)
tx(Inches(0.7),Inches(0.78),Inches(12),Inches(0.5),[[("Help people understand the path. Help projects move through the path.",SERIF,17,SLATE,True,False)],
   [("Leave the path easier for the next person.",SERIF,17,GREEN,True,True)]],anchor=MSO_ANCHOR.TOP,sp=1.05)
# ladder boxes
steps=[("LogicCommons","UNDERSTAND","Free templates, checklists, frameworks","Free",GREEN),
 ("Permit & Bridge","NAVIGATE","Find the path: permits, boards, funding","$250–$750 to start",GOLD),
 ("Stewardship Map","DIAGNOSE","What needs to happen, who owns it, what breaks","$2,500–$7,500",GREEN),
 ("PublicLogic","DELIVER","Sprints · implementation · capacity support","engagement",SLATE)]
bx=Inches(2.55); bw=Inches(8.2); bh=Inches(0.82); y0=2.05; gap=0.235+0.82
for i,(name,verb,desc,price,vc) in enumerate(steps):
    y=Inches(y0+i*gap)
    rect(bx,y,bw,bh,WHITE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,line=CREAMW)
    rect(bx,y,Inches(0.13),bh,vc)  # accent stripe
    tx(bx+Inches(0.4),y,Inches(4.4),bh,[[(name,SERIF,19,SLATE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    tx(bx+Inches(0.4),y+Inches(0.46),Inches(4.6),Inches(0.32),[[(desc,SANS,10.5,MUTE,False,False)]],anchor=MSO_ANCHOR.TOP)
    # verb chip
    tx(bx+Inches(5.1),y+Inches(0.1),Inches(1.9),Inches(0.4),[[(verb,SANS,13,vc,True,False)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    tx(bx+Inches(5.1),y+Inches(0.46),Inches(1.9),Inches(0.32),[[(price,SANS,9.5,GOLD,True,False)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.TOP)
    # down arrow between boxes
    if i<len(steps)-1:
        ar=rect(bx+bw/2-Inches(0.16),y+bh+Inches(0.02),Inches(0.32),Inches(0.2),GOLD,shape=MSO_SHAPE.DOWN_ARROW)
# foundation bar
fy=Inches(y0+len(steps)*gap-0.02)
rect(bx,fy,bw,Inches(0.72),GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tx(bx+Inches(0.4),fy,bw-Inches(0.8),Inches(0.72),[[("CONTINUITY & STEWARDSHIP SYSTEMS   ",SANS,13,WHITE,True,False),("the plumbing, not a product",SERIF,12,RGBColor(0xD8,0xC8,0x9A),False,True)]],anchor=MSO_ANCHOR.MIDDLE)
# footer
tx(Inches(0.7),Inches(7.08),Inches(12),Inches(0.3),[[("Everything is delivered through it.  ",SANS,9,MUTE,False,True),("PublicLogic LLC · Private & Confidential",SANS,9,MUTE,False,False)]],anchor=MSO_ANCHOR.TOP)
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"final_deck","PublicLogic - The Path (one-page visual).pptx")
os.makedirs(os.path.dirname(out),exist_ok=True)
p.save(out); print("wrote",out)
