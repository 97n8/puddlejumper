#!/usr/bin/env python3
"""
PublicLogic — Capabilities & Positioning deck (institutional style), from the
Business Development & Market Positioning Workbook v2.0.
Palette: cream #F5F4F0 / deep slate #1A1D20 / forest green #0F4C3A / ochre accent.
Serif titles (Georgia), sans body (Helvetica Neue). 16:9.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CREAM=RGBColor(0xF5,0xF4,0xF0); SLATE=RGBColor(0x1A,0x1D,0x20); GREEN=RGBColor(0x0F,0x4C,0x3A)
GOLD=RGBColor(0xA9,0x77,0x2F); MUTE=RGBColor(0x6B,0x66,0x60); WHITE=RGBColor(0xFF,0xFF,0xFF)
CREAMW=RGBColor(0xEC,0xEA,0xE3); GREENL=RGBColor(0xE3,0xEA,0xE5); SOFT=RGBColor(0xB9,0xB3,0xA8)
SERIF="Georgia"; SANS="Helvetica Neue"
EW,EH=Inches(13.333),Inches(7.5)

def deck(): p=Presentation(); p.slide_width=EW; p.slide_height=EH; return p
def base(p,bg=CREAM):
    s=p.slides.add_slide(p.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,color,line=None):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=color
    if line: r.line.color.rgb=line; r.line.width=Pt(0.75)
    else: r.line.fill.background()
    r.shadow.inherit=False; return r
def tx(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.06):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,para in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pp.alignment=align; pp.line_spacing=sp
        for (t,fn,sz,col,bold,ital) in para:
            rn=pp.add_run(); rn.text=t; f=rn.font; f.name=fn; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.italic=ital
    return tb
def kick(s,t,col=GOLD): tx(s,Inches(0.85),Inches(0.55),Inches(10),Inches(0.4),[[("   ".join(t.upper()),SANS,11,col,True,False)]])
def rule(s,y=Inches(1.42),w=Inches(0.9),col=GOLD): rect(s,Inches(0.85),y,w,Pt(2.2),col)
def title(s,t,size=30): tx(s,Inches(0.85),Inches(0.95),Inches(11.6),Inches(1.0),[[(t,SERIF,size,SLATE,True,False)]],sp=1.0)
def foot(s): tx(s,Inches(0.85),Inches(7.06),Inches(11.6),Inches(0.3),[[("PUBLICLOGIC LLC   ·   INSTITUTIONAL STEWARDSHIP THROUGH PROJECT DELIVERY",SANS,8,MUTE,False,False)]])
def pg(s,n): tx(s,Inches(12.55),Inches(7.02),Inches(0.6),Inches(0.3),[[(str(n),SANS,9,MUTE,False,False)]],PP_ALIGN.RIGHT)
def section(s,x,y,w,header,bullets,bsize=12.5):
    tx(s,x,y,w,Inches(0.5),[[(header,SANS,13,GREEN,True,False)]])
    tx(s,x,y+Inches(0.5),w,Inches(4.5),[[("—  ",SANS,bsize,GOLD,True,False),(b,SANS,bsize,SLATE,False,False)] for b in bullets],sp=1.14)

p=deck()

# 1 COVER
s=base(p,SLATE)
rect(s,Inches(11.45),0,Inches(1.88),EH,GREEN)
tx(s,Inches(11.45),Inches(3.2),Inches(1.88),Inches(1),[[("steward",SERIF,14,RGBColor(0xCF,0xC9,0xBE),False,True)]],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
kick(s,"Capabilities & Positioning")
rule(s,y=Inches(1.35),w=Inches(1.1))
tx(s,Inches(0.85),Inches(1.95),Inches(10.3),Inches(2.0),[[("PUBLICLOGIC",SERIF,54,CREAM,True,False)],[("Institutional Stewardship through Project Delivery",SERIF,21,GOLD,False,True)]],sp=1.05)
tx(s,Inches(0.85),Inches(4.5),Inches(10.3),Inches(1.2),[
 [("“Implementation is the service.",SERIF,22,CREAM,False,True)],
 [(" Institutional Stewardship is the method.”",SERIF,22,CREAM,False,True)]],sp=1.1)
tx(s,Inches(0.85),Inches(6.3),Inches(10),Inches(0.6),[[("Business Development & Market Positioning  ·  v2.0  ·  June 2026",SANS,11,SOFT,False,False)]])

# 2 THE PROBLEM
s=base(p); kick(s,"The Problem"); rule(s); title(s,"Important work becomes disconnected from the systems that sustain it")
tx(s,Inches(0.85),Inches(2.0),Inches(11.6),Inches(0.6),[[("Organizations don’t fail because they lack ideas. They fail because capacity stops keeping pace with complexity.",SERIF,18,SLATE,False,True)]],sp=1.1)
fails=["Projects stall.","Funding opportunities are missed.","Knowledge leaves.","Responsibilities become unclear.","Plans sit on shelves.","Momentum disappears."]
for i,f in enumerate(fails):
    x=Inches(0.85)+Inches(3.95)*(i%3); y=Inches(3.1)+Inches(0.95)*(i//3)
    rect(s,x,y,Inches(3.7),Inches(0.75),CREAMW); rect(s,x,y,Inches(0.1),Inches(0.75),GOLD)
    tx(s,x+Inches(0.3),y,Inches(3.3),Inches(0.75),[[(f,SANS,13.5,SLATE,True,False)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
tx(s,Inches(0.85),Inches(5.3),Inches(11.6),Inches(0.6),[[("These look like separate problems. They are usually one: institutional capacity is not keeping pace with organizational complexity.",SANS,12.5,GREEN,True,False)]],sp=1.1)
foot(s); pg(s,2)

# 3 THESIS
s=base(p,GREEN); rule(s,y=Inches(2.5),w=Inches(1.1),col=GOLD)
tx(s,Inches(0.85),Inches(1.5),Inches(11),Inches(1),[[("Every organization has a Chair.",SERIF,38,CREAM,True,False)]])
tx(s,Inches(0.85),Inches(2.8),Inches(11),Inches(1.2),[[("Not a seat. A responsibility. A function. A body of knowledge. A place where continuity resides.",SERIF,22,CREAM,False,True)]],sp=1.1)
tx(s,Inches(0.85),Inches(4.6),Inches(11),Inches(1.6),[
 [("When the Chair depends on a single individual, risk rises.",SANS,14,RGBColor(0xDC,0xD7,0xCD),False,False)],
 [("When knowledge is undocumented, projects become fragile.",SANS,14,RGBColor(0xDC,0xD7,0xCD),False,False)],
 [("When turnover occurs, momentum is lost.",SANS,14,RGBColor(0xDC,0xD7,0xCD),False,False)],
 [("PublicLogic exists to reduce those risks.",SANS,14,GOLD,True,False)]],sp=1.4)

# 4 WHAT WE SELL
s=base(p); kick(s,"What We Sell"); rule(s); title(s,"Clients buy outcomes. We deliver them through stewardship.")
rows=[("“Help us figure this out.”","Project Development"),("“Help us pay for this.”","Funding Strategy"),("“Help us get this done.”","Implementation Support")]
y=Inches(2.3)
for i,(need,serv) in enumerate(rows):
    yy=y+Inches(1.05)*i
    rect(s,Inches(0.85),yy,Inches(6.0),Inches(0.85),CREAMW)
    tx(s,Inches(1.15),yy,Inches(5.6),Inches(0.85),[[(need,SERIF,18,SLATE,False,True)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    rect(s,Inches(7.1),yy,Inches(5.35),Inches(0.85),GREEN)
    tx(s,Inches(7.4),yy,Inches(4.9),Inches(0.85),[[(serv,SANS,16,CREAM,True,False)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
tx(s,Inches(0.85),Inches(5.9),Inches(11.6),Inches(0.5),[[("The visible side of PublicLogic. The method underneath is Institutional Stewardship.",SANS,12.5,MUTE,False,True)]])
foot(s); pg(s,4)

# 5 HOW WE DELIVER
s=base(p); kick(s,"How We Deliver"); rule(s); title(s,"Institutional Stewardship")
rect(s,Inches(0.85),Inches(2.0),Inches(11.6),Inches(1.5),CREAMW); rect(s,Inches(0.85),Inches(2.0),Inches(0.12),Inches(1.5),GREEN)
tx(s,Inches(1.2),Inches(2.25),Inches(11.0),Inches(1.1),[[("The deliberate practice of preserving, governing, transferring, and operationalizing organizational knowledge — so critical functions remain durable across turnover, changing priorities, and organizational change.",SERIF,18,SLATE,False,True)]],sp=1.12)
tx(s,Inches(0.85),Inches(4.0),Inches(11.6),Inches(2),[
 [("This is not a slogan. It is an operating discipline.",SANS,14,GREEN,True,False)],
 [("Every PublicLogic engagement applies Institutional Stewardship principles, regardless of project type.",SANS,13,SLATE,False,False)],
 [("Implementation is the service. Institutional Stewardship is the method. Both are true; neither is sacrificed for the other.",SANS,13,SLATE,False,True)]],sp=1.5)
foot(s); pg(s,5)

# 6 SERVICE LINES
s=base(p); kick(s,"Core Service Lines"); rule(s); title(s,"Three services. One discipline.")
svc=[("Project Development","Turn ideas and priorities into actionable projects.",["Project Charter","Scope Definition","Stakeholder Map","Readiness Assessment","Action Roadmap"],"Clarity before resources are committed."),
 ("Funding Strategy","Identify, organize, and pursue sustainable funding.",["Funding Scan","Funding Roadmap","Application Support","Capital Strategy","Partner Coordination"],"Align funding with long-term capacity."),
 ("Implementation Support","Move projects from planning into execution.",["Implementation Framework","Accountability Structure","Progress Tracking","Facilitation","Stakeholder Coordination"],"Work survives beyond its initial champions.")]
for i,(h,purp,dels,obj) in enumerate(svc):
    x=Inches(0.85)+Inches(3.95)*i
    tx(s,x,Inches(1.95),Inches(3.7),Inches(0.4),[[(h,SANS,14,GREEN,True,False)]])
    tx(s,x,Inches(2.4),Inches(3.7),Inches(0.7),[[(purp,SANS,11.5,MUTE,False,True)]],sp=1.05)
    tx(s,x,Inches(3.2),Inches(3.7),Inches(2.3),[[("—  ",SANS,11.5,GOLD,True,False),(d,SANS,11.5,SLATE,False,False)] for d in dels],sp=1.18)
    rect(s,x,Inches(6.05),Inches(3.7),Inches(0.6),GREENL)
    tx(s,x+Inches(0.2),Inches(6.05),Inches(3.4),Inches(0.6),[[(obj,SANS,10.5,GREEN,True,True)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
pg(s,6)

# 7 NATHAN + ALLIE
s=base(p); kick(s,"The Model"); rule(s); title(s,"Institutional systems and human systems are interdependent.")
section(s,Inches(0.85),Inches(2.0),Inches(5.6),"Nathan — Institutional Systems",
 ["Municipal administration, finance, capital planning, procurement.","Grants, economic development, governance, project delivery.","Town Administrator · City Councilor · Town Clerk · Consultant."])
section(s,Inches(6.95),Inches(2.0),Inches(5.5),"Allie — Human Systems",
 ["Organizational readiness, stakeholder engagement, leadership development.","Program design, evaluation, behavioral systems, change management.","Clinical leadership · organizational consulting · workforce systems."])
rect(s,Inches(0.85),Inches(5.6),Inches(11.6),Inches(0.85),GREEN)
tx(s,Inches(1.15),Inches(5.6),Inches(11.0),Inches(0.85),[[("Most firms understand systems OR people. PublicLogic is designed to understand both — because projects fail when either side is ignored.",SANS,13,CREAM,True,True)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
pg(s,7)

# 8 TRANSFER ASSETS
s=base(p); kick(s,"Transfer Assets"); rule(s); title(s,"Every project leaves the client more capable than before.")
assets=[("LogicCommons","Shared frameworks, templates, operational resources."),("Workbooks","Structured implementation and planning guides."),
 ("CaseSpaces","Project-specific operating environments."),("PuddleJumper","Process architecture and workflow infrastructure."),
 ("VAULT","Knowledge preservation and continuity framework.")]
for i,(n,d) in enumerate(assets):
    x=Inches(0.85)+Inches(2.42)*i
    rect(s,x,Inches(2.1),Inches(2.25),Inches(2.4),CREAMW); rect(s,x,Inches(2.1),Inches(2.25),Inches(0.1),GOLD)
    tx(s,x+Inches(0.18),Inches(2.35),Inches(1.9),Inches(0.7),[[(n,SANS,13.5,GREEN,True,False)]])
    tx(s,x+Inches(0.18),Inches(3.15),Inches(1.9),Inches(1.2),[[(d,SANS,10.5,SLATE,False,False)]],sp=1.08)
tx(s,Inches(0.85),Inches(5.0),Inches(11.6),Inches(1),[
 [("The success standard:",SANS,14,GREEN,True,False)],
 [("A successful engagement reduces future dependency. The organization becomes more capable; the process more durable; the knowledge more accessible.",SANS,13,SLATE,False,False)]],sp=1.4)
foot(s); pg(s,8)

# 9 TARGET MARKETS
s=base(p); kick(s,"Target Markets — Tier 1"); rule(s); title(s,"Organizations that already have projects and budgets.")
tm=[("Regional Planning Agencies",["Municipal implementation support, delivery capacity, community engagement.","CMRPC · MRPC · MVPC · Strafford RPC."]),
 ("Engineering Firms",["Public-process support, funding narratives, municipal coordination.","Fuss & O’Neill · Weston & Sampson · Tighe & Bond · BETA."]),
 ("Municipal Consulting Firms",["Capacity, specialized expertise, project support.","Community Paradigm and similar providers — subcontracting paths."]),
 ("Municipalities & the Existing Network",["Temporary capacity, strategic initiatives, capital planning, funding.","Highest trust, shortest sales cycle."])]
for i,(h,b) in enumerate(tm):
    x=Inches(0.85)+Inches(6.0)*(i%2); y=Inches(2.0)+Inches(2.45)*(i//2)
    section(s,x,y,Inches(5.6),h,b)
foot(s); pg(s,9)

# 10 SIGNATURE OFFERINGS
s=base(p); kick(s,"Signature Offerings"); rule(s); title(s,"Three ways to start.")
off=[("Project Development Sprint","Move an idea into an actionable project.","Output: Project Development Roadmap"),
 ("Funding Strategy Sprint","Identify realistic funding pathways and readiness requirements.","Output: Funding Roadmap"),
 ("Implementation Support Partnership","Provide coordination and stewardship capacity through implementation.","Output: Implementation Framework + Accountability Structure")]
for i,(h,p2,o) in enumerate(off):
    x=Inches(0.85)+Inches(3.95)*i
    rect(s,x,Inches(2.1),Inches(3.7),Inches(3.2),CREAMW); rect(s,x,Inches(2.1),Inches(3.7),Inches(0.13),GREEN)
    tx(s,x+Inches(0.3),Inches(2.4),Inches(3.1),Inches(0.9),[[(h,SANS,15,GREEN,True,False)]],sp=1.05)
    tx(s,x+Inches(0.3),Inches(3.4),Inches(3.1),Inches(1.2),[[(p2,SANS,12,SLATE,False,False)]],sp=1.1)
    rect(s,x,Inches(4.75),Inches(3.7),Inches(0.55),GREENL)
    tx(s,x+Inches(0.2),Inches(4.75),Inches(3.4),Inches(0.55),[[(o,SANS,10.5,GREEN,True,True)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
pg(s,10)

# 11 LESSONS LEARNED
s=base(p); kick(s,"Lessons Learned · AI for Impact"); rule(s); title(s,"People buy outcomes.")
rect(s,Inches(0.85),Inches(2.1),Inches(5.6),Inches(2.6),CREAMW)
tx(s,Inches(1.15),Inches(2.3),Inches(5.0),Inches(0.4),[[("WHAT WE THOUGHT THEY WANTED",SANS,11,MUTE,True,False)]])
tx(s,Inches(1.15),Inches(2.8),Inches(5.0),Inches(1.7),[[("AI",SERIF,17,SLATE,False,False)],[("Innovation",SERIF,17,SLATE,False,False)],[("Readiness",SERIF,17,SLATE,False,False)],[("Transformation",SERIF,17,SLATE,False,False)]],sp=1.25)
rect(s,Inches(6.95),Inches(2.1),Inches(5.5),Inches(2.6),GREEN)
tx(s,Inches(7.25),Inches(2.3),Inches(5.0),Inches(0.4),[[("WHAT THEY ACTUALLY WANT",SANS,11,RGBColor(0xCF,0xC9,0xBE),True,False)]])
tx(s,Inches(7.25),Inches(2.8),Inches(5.0),Inches(1.7),[[("Capacity",SERIF,17,CREAM,True,False)],[("Funding",SERIF,17,CREAM,True,False)],[("Progress",SERIF,17,CREAM,True,False)],[("Execution",SERIF,17,CREAM,True,False)]],sp=1.25)
tx(s,Inches(0.85),Inches(5.1),Inches(11.6),Inches(0.9),[[("The methodology matters. But outcomes drive the purchasing decision. Lead with the known problem — not a hidden one.",SANS,13.5,GOLD,True,True)]],sp=1.1)
foot(s); pg(s,11)

# 12 PROMISE / POSITIONING
s=base(p,SLATE); rule(s,y=Inches(2.4),w=Inches(1.1),col=GOLD)
kick(s,"The Promise")
tx(s,Inches(0.85),Inches(1.5),Inches(11.5),Inches(1.5),[[("We move important projects from planning to implementation —",SERIF,28,CREAM,True,False)],[("and keep the knowledge, systems, and structures that sustain them intact.",SERIF,28,GOLD,False,True)]],sp=1.08)
tx(s,Inches(0.85),Inches(4.0),Inches(11.5),Inches(1.4),[[("PublicLogic helps communities and mission-driven organizations move important projects from planning to implementation. We do this through Institutional Stewardship — preserving, governing, transferring, and operationalizing the knowledge required for long-term success.",SANS,13.5,RGBColor(0xDC,0xD7,0xCD),False,False)]],sp=1.2)
rect(s,Inches(0.85),Inches(5.9),Inches(11.6),Inches(0.85),GREEN)
tx(s,Inches(1.15),Inches(5.9),Inches(11.0),Inches(0.85),[[("The market buys project delivery. PublicLogic delivers Institutional Stewardship. Both statements are true; neither is sacrificed for the other.",SANS,13,CREAM,True,True)]],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)

out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"final_deck","PublicLogic - Capabilities & Positioning (2026-06).pptx")
os.makedirs(os.path.dirname(out),exist_ok=True)
p.save(out); print("wrote",out,"| slides:",len(p.slides))
